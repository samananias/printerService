"""
spike_t6_office.py — Office Printing Spike (docs/MULTI_FORMAT_PLAN.md §14, T6)

Run this ON the print-server PC, from the project root, AFTER installing
LibreOffice (https://www.libreoffice.org, or: winget install TheDocumentFoundation.LibreOffice):

    .venv\\Scripts\\pip install python-docx openpyxl python-pptx
    .venv\\Scripts\\python spike_t6_office.py

Generates three REAL office documents, converts each with the service's
real OfficeProcessor (LibreOffice headless — the exact production path,
including its timeout and fresh-profile handling), prints the PDFs via
SumatraPDF, and times every conversion:

  1. DOCX — a table-heavy document (a 12x4 bordered table + headings)
  2. XLSX — a spreadsheet with a defined PRINT AREA + landscape page setup
  3. PPTX — a 16:9 deck (2 slides, landscape)

PASS criteria — judge the PAPER (the script cannot see it):
  [ ] DOCX: table fits the page width, borders visible, no cut columns
  [ ] XLSX: ONLY the print area prints, in landscape, on one page
  [ ] PPTX: slides fill the page in landscape (16:9)
Also note each conversion time — if the machine takes > 30 s per document,
consider raising CONVERT_TIMEOUT_S in .env.
Record the summary in SOURCE_OF_TRUTH Section 5, like the T4/T5 entries.
These results are the Phase 3 acceptance gate.
"""

import argparse
import shutil
import subprocess
import sys
import tempfile
import time
from pathlib import Path

LINE = "=" * 64


def banner(text: str) -> None:
    print("\n" + LINE)
    print(text)
    print(LINE)


def find_printer() -> str:
    """Prefer the L3210 by name, fall back to the Windows default."""
    import win32print

    flags = win32print.PRINTER_ENUM_LOCAL | win32print.PRINTER_ENUM_CONNECTIONS
    names = sorted(p[2] for p in win32print.EnumPrinters(flags))
    if not names:
        raise RuntimeError("No printers found — is the Epson installed on this PC?")
    for name in names:
        if "L3210" in name:
            return name
    return names[0]


def make_test_documents(folder: Path) -> list[tuple[str, Path]]:
    """Generate the three spike documents with the python-* helper libs."""
    out: list[tuple[str, Path]] = []

    # 1. DOCX — table-heavy (the documented DOCX quality check).
    import docx

    doc = docx.Document()
    doc.add_heading("T6 office spike — table layout check", 1)
    doc.add_paragraph(
        "If this table fits the page width with visible borders and no cut "
        "columns, DOCX table conversion works on this machine."
    )
    table = doc.add_table(rows=12, cols=4)
    table.style = "Table Grid"
    for row in range(12):
        for col in range(4):
            table.cell(row, col).text = f"r{row + 1}c{col + 1} — some cell content"
    path = folder / "t6_1_table.docx"
    doc.save(path)
    out.append(("1 DOCX table-heavy", path))

    # 2. XLSX — print area + landscape (the documented XLSX quality check:
    #    ONLY the print area may come out of the printer, in landscape).
    import openpyxl

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Spike"
    for row in range(1, 21):
        for col in range(1, 5):
            sheet.cell(row=row, column=col, value=f"cell r{row}c{col}")
    sheet.cell(row=1, column=6, value="OUTSIDE the print area — must NOT print")
    sheet.print_area = "A1:D20"
    sheet.page_setup.orientation = "landscape"
    path = folder / "t6_2_printarea.xlsx"
    workbook.save(path)
    out.append(("2 XLSX with print area", path))

    # 3. PPTX — 16:9 deck (slide size should drive a landscape PDF page).
    from pptx import Presentation
    from pptx.util import Inches

    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    for number in (1, 2):
        slide = deck.slides.add_slide(deck.slide_layouts[6])  # blank layout
        box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
        box.text_frame.text = f"T6 office spike — slide {number} of 2 (16:9)"
    path = folder / "t6_3_deck.pptx"
    deck.save(path)
    out.append(("3 PPTX 16:9 deck", path))

    return out


def print_pdf(sumatra: str, pdf_path: Path, printer_name: str) -> None:
    """The service's exact print invocation (see app/printer/windows.py)."""
    result = subprocess.run(
        [sumatra, "-print-to", printer_name, "-silent", str(pdf_path)],
        capture_output=True,
        timeout=180,
    )
    if result.returncode != 0:
        raise RuntimeError(
            f"SumatraPDF exited with code {result.returncode}: "
            f"{result.stderr.decode(errors='replace').strip()}"
        )


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.parse_args()

    banner("T6 OFFICE SPIKE — run this ON the PC the printer is plugged into")

    try:
        from app.printer.windows import find_sumatra
        from app.processors.office import OfficeProcessor, find_soffice
    except ImportError as exc:
        print(f"Cannot import the app ({exc}). Run from the project root:")
        print("    .venv\\Scripts\\python spike_t6_office.py")
        return 1

    try:
        import win32print  # noqa: F401  (pywin32 presence check, like T1)
    except ImportError:
        print("pywin32 is not installed here:  pip install pywin32")
        return 1

    processor = OfficeProcessor()
    if not processor.available():
        print(
            "Office conversion is NOT available on this machine:\n"
            "  - install LibreOffice (winget install TheDocumentFoundation.LibreOffice)\n"
            "  - and/or set ENABLE_OFFICE=1 / LO_PATH in .env\n"
            f"  find_soffice() -> {find_soffice()!r}, ENABLE_OFFICE from .env"
        )
        return 1

    printer_name = find_printer()
    sumatra = find_sumatra()
    if not sumatra:
        print("SumatraPDF not found — install it or set SUMATRA_PATH in .env")
        return 1

    print(f"\nPrinter:    {printer_name}")
    print(f"Sumatra:    {sumatra}")
    print(f"LibreOffice: {find_soffice()}")
    print("\n>>> Keep paper loaded and watch the physical printer.")
    input("Press Enter when ready...")

    temp_dir = Path(tempfile.mkdtemp(prefix="spike_t6_"))
    results: list[tuple[str, str, str]] = []
    try:
        for name, doc_path in make_test_documents(temp_dir):
            try:
                started = time.perf_counter()
                pdf_path = processor.process(doc_path, temp_dir)
                seconds = time.perf_counter() - started
                print_pdf(sumatra, pdf_path, printer_name)
                results.append(
                    (
                        f"T6 {name}",
                        "PASS",
                        f"converted in {seconds:.1f}s, print accepted — CHECK PAPER",
                    )
                )
            except Exception as exc:
                results.append((f"T6 {name}", "FAIL", str(exc)))
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)

    banner("SUMMARY")
    for name, status, detail in results:
        print(f"[{status:4}] {name}: {detail}")

    print(
        "\nNow judge the paper:\n"
        "  [ ] DOCX table fits the width, borders visible, no cut columns\n"
        "  [ ] XLSX: ONLY A1:D20 printed (no 'OUTSIDE' cell!), landscape\n"
        "  [ ] PPTX slides fill the page, landscape (16:9)\n"
        "\nRecord the results in SOURCE_OF_TRUTH Section 5 (like the T4/T5\n"
        "entries) — they are the Phase 3 acceptance gate."
    )
    return 0 if all(r[1] == "PASS" for r in results) else 2


if __name__ == "__main__":
    sys.exit(main())
